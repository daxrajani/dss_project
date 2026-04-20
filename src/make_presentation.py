#!/usr/bin/env python3
"""
Generate DSS Final Presentation — COEN 6731 Winter 2026
Authors: Dax Rajani (40294118) & Harsh Ahuja (40301748)

8 slides, widescreen 16:9 (13.33" × 7.5")
  1. Title
  2. Problem Statement
  3. System Architecture        ← mandatory diagram
  4. Implementation Structure   ← mandatory diagram
  5. Pipeline: 5 Stages
  6. Performance Optimizations
  7. Scaling Benchmark Results
  8. Lessons Learned

Usage:
    source dss_env/bin/activate
    python src/make_presentation.py
Output:
    DSS_Final_Presentation.pptx
"""

import os
import sys
from io import BytesIO

# ── matplotlib (for charts) ──────────────────────────────────────────────────
try:
    import matplotlib
    matplotlib.use("Agg")
    import matplotlib.pyplot as plt
    import numpy as np
except ImportError:
    print("ERROR: pip install matplotlib numpy"); sys.exit(1)

# ── python-pptx ──────────────────────────────────────────────────────────────
try:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
except ImportError:
    print("ERROR: pip install python-pptx"); sys.exit(1)

# ── Paths ─────────────────────────────────────────────────────────────────────
SCRIPT_DIR  = os.path.dirname(os.path.abspath(__file__))
PROJECT_DIR = os.path.dirname(SCRIPT_DIR)
CHARTS_DIR  = os.path.join(PROJECT_DIR, "results", "charts")
OUT_PATH    = os.path.join(PROJECT_DIR, "DSS_Final_Presentation.pptx")

# ── Slide dimensions ─────────────────────────────────────────────────────────
W = Inches(13.33)   # 16:9 widescreen width
H = Inches(7.50)    # 16:9 widescreen height

# ── Color palette ─────────────────────────────────────────────────────────────
NAVY     = RGBColor(0x0D, 0x1B, 0x3E)
BLUE     = RGBColor(0x1A, 0x5E, 0xA8)
LTBLUE   = RGBColor(0x5B, 0x9B, 0xD5)
TEAL     = RGBColor(0x0E, 0x7A, 0x8A)
GREEN    = RGBColor(0x1E, 0x8A, 0x44)
ORANGE   = RGBColor(0xE0, 0x72, 0x10)
RED      = RGBColor(0xC0, 0x2B, 0x2B)
PURPLE   = RGBColor(0x6B, 0x27, 0x8B)
GOLD     = RGBColor(0xF4, 0xA4, 0x1B)
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
OFFWHITE = RGBColor(0xF4, 0xF6, 0xFA)
LGRAY    = RGBColor(0xCC, 0xD5, 0xE2)
GRAY     = RGBColor(0x55, 0x65, 0x7A)
DARK     = RGBColor(0x1A, 0x1A, 0x2E)


# ─────────────────────────────────────────────────────────────────────────────
#  LOW-LEVEL HELPERS
# ─────────────────────────────────────────────────────────────────────────────

def new_prs() -> Presentation:
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H
    return prs


def blank_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])  # blank layout


def set_bg(slide, rgb: RGBColor):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = rgb


def box(slide, l, t, w, h, fill=None, line=None, lw=Pt(1.5), rounded=False):
    """Add a rectangle (or rounded rectangle) shape."""
    stype = (MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE if rounded
             else MSO_AUTO_SHAPE_TYPE.RECTANGLE)
    s = slide.shapes.add_shape(stype, l, t, w, h)
    if fill:
        s.fill.solid()
        s.fill.fore_color.rgb = fill
    else:
        s.fill.background()
    if line:
        s.line.color.rgb = line
        s.line.width = lw
    else:
        s.line.fill.background()
    return s


def txt(slide, text, l, t, w, h,
        size=Pt(14), bold=False, color=WHITE,
        align=PP_ALIGN.LEFT, italic=False, wrap=True):
    """Add a text box."""
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = size
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.italic = italic
    return tb


def header(slide, title: str, badge: str = ""):
    """Dark header bar with title."""
    box(slide, 0, 0, W, Inches(0.68), fill=NAVY)
    txt(slide, title, Inches(0.35), Inches(0.07), Inches(11.5), Inches(0.58),
        size=Pt(25), bold=True, color=WHITE, align=PP_ALIGN.LEFT)
    if badge:
        txt(slide, badge, Inches(11.5), Inches(0.14), Inches(1.7), Inches(0.42),
            size=Pt(12), color=LTBLUE, align=PP_ALIGN.RIGHT, italic=True)


def img(slide, path_or_buf, l, t, w, h):
    """Embed an image from file path or BytesIO."""
    slide.shapes.add_picture(path_or_buf, l, t, w, h)


# ─────────────────────────────────────────────────────────────────────────────
#  CHART GENERATORS
# ─────────────────────────────────────────────────────────────────────────────

_MPL_STYLE = {
    "figure.facecolor": "#F4F6FA",
    "axes.facecolor":   "#FFFFFF",
    "axes.spines.top":  False,
    "axes.spines.right":False,
    "font.family":      "sans-serif",
}


def _buf(fig) -> BytesIO:
    b = BytesIO()
    fig.savefig(b, format="png", dpi=160, bbox_inches="tight")
    plt.close(fig)
    b.seek(0)
    return b


def chart_scaling_time() -> BytesIO:
    """Grouped bar chart: execution time vs workers × data size."""
    with plt.rc_context(_MPL_STYLE):
        fig, ax = plt.subplots(figsize=(7.5, 4.2))
        workers = [2, 4, 5]
        data = {
            "1 GB":  [673.6,  282.4, 256.9],
            "5 GB":  [1021.8, 569.0, 409.7],
            "10 GB": [1378.9, 799.8, 674.0],
        }
        colors = ["#1A5EA8", "#E07210", "#C02B2B"]
        x = np.arange(len(workers))
        w = 0.26
        for i, (label, vals) in enumerate(data.items()):
            bars = ax.bar(x + i*w, vals, w, label=label,
                          color=colors[i], edgecolor="white", linewidth=1.2)
            ax.bar_label(bars, labels=[f"{v:.0f}s" for v in vals],
                         padding=3, fontsize=8.5, fontweight="bold", color=colors[i])
        ax.set_xticks(x + w)
        ax.set_xticklabels([f"{w}w" for w in workers], fontsize=11)
        ax.set_ylabel("Total Pipeline Time (seconds)", fontsize=11)
        ax.set_title("Execution Time  vs.  Cluster Size", fontsize=13, fontweight="bold", pad=8)
        ax.legend(title="Data Size", fontsize=10, title_fontsize=10)
        ax.set_ylim(0, 1650)
        ax.grid(axis="y", alpha=0.35)
        fig.tight_layout()
        return _buf(fig)


def chart_speedup() -> BytesIO:
    """Line chart: actual speedup vs ideal."""
    with plt.rc_context(_MPL_STYLE):
        fig, ax = plt.subplots(figsize=(5.5, 4.2))
        workers = [2, 4, 5]
        ideal   = [1.0, 2.0, 2.5]
        s1  = [1.0, 673.6/282.4, 673.6/256.9]   # 1 GB
        s5  = [1.0, 1021.8/569.0, 1021.8/409.7]  # 5 GB
        s10 = [1.0, 1378.9/799.8, 1378.9/674.0]  # 10 GB
        ax.plot(workers, ideal, "k--", linewidth=1.8, alpha=0.5, label="Ideal")
        ax.plot(workers, s1,  "o-", color="#1A5EA8", lw=2.5, ms=8, label=f"1 GB")
        ax.plot(workers, s5,  "s-", color="#E07210", lw=2.5, ms=8, label=f"5 GB")
        ax.plot(workers, s10, "^-", color="#C02B2B", lw=2.5, ms=8, label=f"10 GB")
        # annotate 5-worker points
        for val, col in [(s1[-1], "#1A5EA8"), (s5[-1], "#E07210"), (s10[-1], "#C02B2B")]:
            ax.annotate(f"{val:.2f}×", xy=(5, val), xytext=(4.62, val + 0.07),
                        fontsize=10, color=col, fontweight="bold")
        ax.set_xticks(workers)
        ax.set_xticklabels([f"{w}w" for w in workers], fontsize=11)
        ax.set_ylabel("Speedup over 2-worker baseline", fontsize=11)
        ax.set_title("Speedup  vs.  Cluster Size", fontsize=13, fontweight="bold", pad=8)
        ax.legend(fontsize=10)
        ax.set_ylim(0.5, 3.0)
        ax.grid(alpha=0.35)
        fig.tight_layout()
        return _buf(fig)


def chart_optimizations() -> BytesIO:
    """3-panel chart: broadcast join, partitioning, skew."""
    with plt.rc_context(_MPL_STYLE):
        fig, axes = plt.subplots(1, 3, figsize=(13, 4.0))

        # Panel 1 — Broadcast Join
        ax = axes[0]
        methods = ["Shuffle Join\n(Baseline)", "Broadcast Join"]
        times   = [11.33, 7.15]
        clrs    = ["#1A5EA8", "#1E8A44"]
        bars = ax.bar(methods, times, color=clrs, width=0.45,
                      edgecolor="white", linewidth=1.5)
        ax.bar_label(bars, labels=[f"{t}s" for t in times],
                     padding=4, fontsize=13, fontweight="bold")
        ax.annotate("1.58× faster", xy=(1, 7.15), xytext=(0.5, 10.2),
                    fontsize=12, fontweight="bold", color="#1E8A44",
                    arrowprops=dict(arrowstyle="->", color="#1E8A44", lw=1.8))
        ax.set_ylim(0, 14)
        ax.set_title("Broadcast Join\nOptimization", fontsize=12, fontweight="bold")
        ax.set_ylabel("Time (seconds)", fontsize=10)

        # Panel 2 — Partition Pruning
        ax2 = axes[1]
        qtype   = ["Single-day\nQuery", "7-day Range\nQuery"]
        speedup = [3.11, 1.73]
        clrs2   = ["#E07210", "#E07210"]
        bars2 = ax2.bar(qtype, speedup, color=clrs2, width=0.45,
                        edgecolor="white", linewidth=1.5)
        ax2.bar_label(bars2, labels=[f"{s}×" for s in speedup],
                      padding=4, fontsize=13, fontweight="bold", color="#E07210")
        ax2.axhline(y=1.0, color="#556478", linestyle="--", lw=1.5, label="Baseline (flat)")
        ax2.set_ylim(0, 4.2)
        ax2.set_title("Date-Partitioned Parquet\n(Partition Pruning)", fontsize=12, fontweight="bold")
        ax2.set_ylabel("Speedup over flat Parquet", fontsize=10)
        ax2.legend(fontsize=9)

        # Panel 3 — Skew Analysis
        ax3 = axes[2]
        strats = ["No Salting\n(Baseline)", "With Salting\n(20 buckets)"]
        times3 = [34.6, 79.1]
        clrs3  = ["#1E8A44", "#C02B2B"]
        bars3 = ax3.bar(strats, times3, color=clrs3, width=0.45,
                        edgecolor="white", linewidth=1.5)
        ax3.bar_label(bars3, labels=[f"{t}s" for t in times3],
                      padding=4, fontsize=13, fontweight="bold")
        ax3.set_ylim(0, 105)
        ax3.set_title("Skew Mitigation Analysis\n(Attribution Join)", fontsize=12, fontweight="bold")
        ax3.set_ylabel("Time (seconds)", fontsize=10)
        ax3.text(0.5, 88, "Top-1 user = 0.02%\nLow skew → salting hurt",
                 ha="center", fontsize=10, color="#556478", style="italic",
                 bbox=dict(boxstyle="round,pad=0.4", fc="#FFF3CD", ec="#E0C000"))

        fig.tight_layout(pad=2.2)
        return _buf(fig)


def chart_throughput() -> BytesIO:
    """Throughput scaling line chart."""
    with plt.rc_context(_MPL_STYLE):
        fig, ax = plt.subplots(figsize=(7.5, 4.2))
        workers = [2, 4, 5]
        tp = {
            "1 GB":  [12607,  30070,  33055],
            "5 GB":  [41543,  74602,  103609],
            "10 GB": [61569, 106148, 125961],
        }
        markers = ["o", "s", "^"]
        colors  = ["#1A5EA8", "#E07210", "#C02B2B"]
        for (label, vals), m, c in zip(tp.items(), markers, colors):
            ax.plot(workers, [v/1000 for v in vals], f"{m}-",
                    color=c, lw=2.5, ms=9, label=label)
            ax.annotate(f"{vals[-1]/1000:.0f}K",
                        xy=(5, vals[-1]/1000), xytext=(5.08, vals[-1]/1000),
                        fontsize=9.5, color=c, fontweight="bold", va="center")
        ax.set_xticks(workers)
        ax.set_xticklabels([f"{w} workers" for w in workers], fontsize=11)
        ax.set_ylabel("Throughput  (K rows / sec)", fontsize=11)
        ax.set_title("Throughput  vs.  Cluster Size", fontsize=13, fontweight="bold", pad=8)
        ax.legend(title="Data Size", fontsize=10, title_fontsize=10)
        ax.grid(alpha=0.35)
        ax.set_xlim(1.7, 5.8)
        fig.tight_layout()
        return _buf(fig)


# ─────────────────────────────────────────────────────────────────────────────
#  SLIDE BUILDERS
# ─────────────────────────────────────────────────────────────────────────────

def slide_title(prs):
    """Slide 1 — Title."""
    slide = blank_slide(prs)
    set_bg(slide, NAVY)

    # Gold accent stripe
    box(slide, 0, 0, Inches(0.22), H, fill=GOLD)

    # Title
    txt(slide, "Scalable E-Commerce Analytics Pipeline",
        Inches(0.52), Inches(1.4), Inches(12.5), Inches(1.3),
        size=Pt(42), bold=True, color=WHITE, align=PP_ALIGN.LEFT)

    # Subtitle
    txt(slide, "PySpark on GCP Dataproc  ·  42M+ Events  ·  5-Stage Distributed Pipeline",
        Inches(0.52), Inches(2.9), Inches(12.5), Inches(0.62),
        size=Pt(21), color=LTBLUE, align=PP_ALIGN.LEFT)

    # Divider
    box(slide, Inches(0.52), Inches(3.68), Inches(11.8), Inches(0.045), fill=GOLD)

    # Authors / course
    txt(slide, "Dax Rajani (40294118)  ·  Harsh Ahuja (40301748)",
        Inches(0.52), Inches(3.88), Inches(11), Inches(0.52),
        size=Pt(18), bold=True, color=WHITE, align=PP_ALIGN.LEFT)
    txt(slide, "COEN 6731 — Distributed Software Systems  ·  Concordia University, Winter 2026",
        Inches(0.52), Inches(4.46), Inches(12.5), Inches(0.48),
        size=Pt(14), color=LGRAY, align=PP_ALIGN.LEFT)

    # Stat badges
    stats = [
        ("42.4 M", "events"),
        ("742 K",  "orders"),
        ("166 K",  "products"),
        ("5",      "pipeline stages"),
        ("9",      "benchmark runs"),
    ]
    bw, bh = Inches(2.28), Inches(0.88)
    for i, (val, lbl) in enumerate(stats):
        bx = Inches(0.52) + i * (bw + Inches(0.12))
        by = Inches(6.28)
        b = box(slide, bx, by, bw, bh, fill=BLUE, rounded=True)
        txt(slide, val, bx, by + Inches(0.04), bw, Inches(0.45),
            size=Pt(23), bold=True, color=GOLD, align=PP_ALIGN.CENTER)
        txt(slide, lbl, bx, by + Inches(0.47), bw, Inches(0.34),
            size=Pt(12), color=LGRAY, align=PP_ALIGN.CENTER)
    return slide


def slide_problem(prs):
    """Slide 2 — Problem Statement."""
    slide = blank_slide(prs)
    set_bg(slide, OFFWHITE)
    header(slide, "Problem Statement", "Slide 1 / 7")

    # ── Left panel — challenges ──────────────────────────────────────────
    box(slide, Inches(0.28), Inches(0.78), Inches(6.05), Inches(6.0),
        fill=WHITE, line=LGRAY, lw=Pt(1))
    txt(slide, "The Challenge", Inches(0.45), Inches(0.90), Inches(5.7), Inches(0.5),
        size=Pt(17), bold=True, color=NAVY)
    box(slide, Inches(0.44), Inches(1.45), Inches(0.06), Inches(0.26), fill=GOLD)

    challenges = [
        ("Scale",
         "42M+ events: pandas times out. Need distributed Spark to process in minutes."),
        ("Missing Features",
         "Raw data lacks device & referrer columns — synthesized via deterministic CRC32 hash on user_session."),
        ("Multi-Stage Logic",
         "Sessionization → Funnel → Attribution → Anomaly must form a reproducible, timed pipeline."),
        ("Cost vs. Performance",
         "Is adding more workers worth the GCP cost? Quantify actual speedup vs. ideal scaling."),
    ]
    cy = Inches(1.6)
    for title, desc in challenges:
        box(slide, Inches(0.44), cy + Inches(0.11), Inches(0.13), Inches(0.13),
            fill=GOLD, rounded=True)
        txt(slide, title + ":", Inches(0.65), cy, Inches(5.45), Inches(0.38),
            size=Pt(13), bold=True, color=NAVY)
        txt(slide, desc, Inches(0.65), cy + Inches(0.36), Inches(5.45), Inches(0.7),
            size=Pt(11.5), color=GRAY, wrap=True)
        cy += Inches(1.24)

    # ── Right panel — dataset + questions ───────────────────────────────
    box(slide, Inches(6.62), Inches(0.78), Inches(6.4), Inches(6.0),
        fill=WHITE, line=LGRAY, lw=Pt(1))
    txt(slide, "Dataset & Research Questions", Inches(6.8), Inches(0.90), Inches(6.0), Inches(0.5),
        size=Pt(17), bold=True, color=NAVY)
    box(slide, Inches(6.8), Inches(1.45), Inches(0.06), Inches(0.26), fill=BLUE)

    # Dataset card
    ds = box(slide, Inches(6.8), Inches(1.68), Inches(6.0), Inches(1.62),
             fill=NAVY, rounded=True)
    txt(slide, "REES46 eCommerce Behavior Data  (Kaggle, Oct 2019)",
        Inches(6.92), Inches(1.74), Inches(5.78), Inches(0.42),
        size=Pt(12), bold=True, color=GOLD)
    txt(slide,
        "  42,448,764 events   ·   742,849 orders   ·   166,794 products\n"
        "  Columns: event_time · event_type · product_id · category_code\n"
        "            brand · price · user_id · user_session\n"
        "  Synthesized: device (mobile / desktop / tablet)  ·  referrer channel",
        Inches(6.92), Inches(2.18), Inches(5.78), Inches(1.05),
        size=Pt(10.5), color=OFFWHITE, wrap=True)

    questions = [
        ("Q1", "What is the view → cart → purchase conversion funnel by category & channel?"),
        ("Q2", "Which referrer channel drives the most revenue? (last-touch attribution)"),
        ("Q3", "Are there anomalous SPIKE / DROP days in daily purchases? (z-score)"),
        ("Q4", "How does throughput scale with cluster size and data volume?"),
    ]
    qy = Inches(3.45)
    for qnum, qtext in questions:
        qb = box(slide, Inches(6.8), qy, Inches(0.52), Inches(0.38), fill=BLUE, rounded=True)
        txt(slide, qnum, Inches(6.8), qy, Inches(0.52), Inches(0.38),
            size=Pt(11), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        txt(slide, qtext, Inches(7.42), qy, Inches(5.4), Inches(0.38),
            size=Pt(11.5), color=NAVY)
        qy += Inches(0.62)
    return slide


def slide_architecture(prs):
    """Slide 3 — System Architecture (mandatory figure)."""
    slide = blank_slide(prs)
    set_bg(slide, OFFWHITE)
    header(slide, "System Architecture", "Slide 2 / 7")

    # GCP Dataproc outer border
    box(slide, Inches(0.28), Inches(0.76), Inches(12.76), Inches(5.58),
        line=BLUE, lw=Pt(2.2))
    txt(slide, "GCP Dataproc Cluster  (1 Master + 2 / 4 / 5 Workers, n1-standard-4, us-central1)",
        Inches(0.46), Inches(0.76), Inches(11), Inches(0.44),
        size=Pt(12), bold=True, color=BLUE)

    # ── Pipeline stage flow ──────────────────────────────────────────────
    SX   = Inches(2.28)    # x of first stage
    SY   = Inches(1.44)    # y of stage top
    SW   = Inches(1.82)    # stage box width
    SH   = Inches(1.32)    # stage box height
    SGAP = Inches(0.24)    # gap between stages

    stages = [
        ("1\nETL &\nEnrichment",    "CRC32 → device\n& referrer\nParquet write",    GREEN),
        ("2\nSession-\nization",    "30-min inactivity\nlag() window\nsession_id",  BLUE),
        ("3\nFunnel\nAnalysis",     "view→cart→buy\nby category\ndevice/referrer", ORANGE),
        ("4\nAttribution\nLast-Touch","24-hr join\nrow_number()\nnon-direct ref",  RED),
        ("5\nAnomaly\nDetection",   "7-day z-score\n|z|>2 → SPIKE\nper-category", PURPLE),
    ]

    # GCS input
    gin = box(slide, Inches(0.36), Inches(1.7), Inches(1.62), Inches(0.96),
              fill=TEAL, rounded=True)
    txt(slide, "GCS Input\nCSV\n(42M rows)", Inches(0.36), Inches(1.72),
        Inches(1.62), Inches(0.96), size=Pt(10.5), bold=True,
        color=WHITE, align=PP_ALIGN.CENTER)
    # arrow GCS→stage1
    box(slide, Inches(1.98), Inches(2.16), Inches(0.32), Inches(0.04), fill=GRAY)

    for i, (name, detail, color) in enumerate(stages):
        bx = SX + i * (SW + SGAP)
        b = box(slide, bx, SY, SW, SH, fill=color, rounded=True)
        # stage label (multiline)
        txt(slide, name, bx, SY + Inches(0.05), SW, Inches(0.62),
            size=Pt(11.5), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        txt(slide, detail, bx, SY + Inches(0.66), SW, Inches(0.62),
            size=Pt(9.5), color=OFFWHITE, align=PP_ALIGN.CENTER)
        # arrow to next
        if i < len(stages) - 1:
            ax2 = bx + SW
            box(slide, ax2, SY + SH/2 - Inches(0.02),
                SGAP, Inches(0.04), fill=GRAY)

    # GCS output
    gout_x = SX + len(stages) * (SW + SGAP)
    box(slide, gout_x, Inches(1.7), Inches(1.45), Inches(0.96),
        fill=TEAL, rounded=True)
    box(slide, gout_x - SGAP, SY + SH/2 - Inches(0.02), SGAP, Inches(0.04), fill=GRAY)
    txt(slide, "GCS Output\nParquet\n(results)", gout_x, Inches(1.72),
        Inches(1.45), Inches(0.96), size=Pt(10.5), bold=True,
        color=WHITE, align=PP_ALIGN.CENTER)

    # ── Stage timing labels ──────────────────────────────────────────────
    timings = ["224s", "185s", "incl.", "64s", "16s"]
    for i, t in enumerate(timings):
        bx = SX + i * (SW + SGAP)
        txt(slide, t, bx, SY + SH + Inches(0.06), SW, Inches(0.3),
            size=Pt(9.5), color=GRAY, align=PP_ALIGN.CENTER, italic=True)
    txt(slide, "5w / 10GB", Inches(0.36), SY + SH + Inches(0.06), Inches(1.62), Inches(0.3),
        size=Pt(9), color=GRAY, align=PP_ALIGN.CENTER, italic=True)

    # ── Worker nodes diagram ─────────────────────────────────────────────
    wn_y = Inches(3.3)
    box(slide, Inches(0.36), wn_y, Inches(12.7), Inches(2.74),
        line=LGRAY, lw=Pt(1.2))
    txt(slide, "Cluster Topology", Inches(0.52), wn_y + Inches(0.06), Inches(4), Inches(0.38),
        size=Pt(12), bold=True, color=NAVY)

    # Master node
    mn = box(slide, Inches(0.52), wn_y + Inches(0.5), Inches(2.0), Inches(0.85),
             fill=NAVY, rounded=True)
    txt(slide, "Master Node\nSparkContext / YARN RM", Inches(0.52), wn_y + Inches(0.5),
        Inches(2.0), Inches(0.85), size=Pt(10), bold=True,
        color=WHITE, align=PP_ALIGN.CENTER)

    # Worker nodes
    worker_labels = ["Worker 1\nn1-std-4\n4 vCPU / 15GB", "Worker 2\nn1-std-4\n4 vCPU / 15GB",
                     "Worker 3\nn1-std-4\n4 vCPU / 15GB", "Worker 4\nn1-std-4\n4 vCPU / 15GB",
                     "Worker 5\nn1-std-4\n4 vCPU / 15GB"]
    ww, wh = Inches(2.0), Inches(0.85)
    for j, lbl in enumerate(worker_labels):
        wx = Inches(3.0) + j * (ww + Inches(0.2))
        wb = box(slide, wx, wn_y + Inches(0.5), ww, wh, fill=BLUE, rounded=True)
        txt(slide, lbl, wx, wn_y + Inches(0.5), ww, wh,
            size=Pt(9.5), color=WHITE, align=PP_ALIGN.CENTER)
        # arrow from master
        box(slide, Inches(2.52), wn_y + Inches(0.88),
            wx - Inches(2.52), Inches(0.04), fill=LGRAY)

    # Arrow from workers to GCS
    txt(slide, "↕  GCS (gs://dss-project-dax)  —  Parquet read/write",
        Inches(3.0), wn_y + Inches(1.56), Inches(9.5), Inches(0.35),
        size=Pt(11), color=TEAL, bold=True, align=PP_ALIGN.CENTER)

    # ── Legend / footnote ────────────────────────────────────────────────
    legend = [
        ("GCS", TEAL), ("ETL", GREEN), ("Sessionize", BLUE),
        ("Funnel", ORANGE), ("Attribution", RED), ("Anomaly", PURPLE),
    ]
    lx = Inches(0.36)
    for (lbl, c) in legend:
        box(slide, lx, Inches(6.48), Inches(0.24), Inches(0.24), fill=c, rounded=True)
        txt(slide, lbl, lx + Inches(0.3), Inches(6.46), Inches(1.4), Inches(0.3),
            size=Pt(10), color=NAVY)
        lx += Inches(1.8)
    return slide


def slide_implementation(prs):
    """Slide 4 — Implementation Structure (mandatory figure)."""
    slide = blank_slide(prs)
    set_bg(slide, OFFWHITE)
    header(slide, "Implementation Structure", "Slide 3 / 7")

    # 4 columns: Data Layer | Core Pipeline | Optimization | Cloud/Infra
    COL_Y    = Inches(1.0)
    COL_H_LBL= Inches(0.4)
    ROW_H    = Inches(0.82)
    ROW_GAP  = Inches(0.16)

    columns = [
        {
            "label": "DATA LAYER",
            "x": Inches(0.28),
            "w": Inches(2.8),
            "color": TEAL,
            "items": [
                ("REES46 CSV\n42M rows / 5.3 GB", TEAL),
                ("enriched_events/\nParquet (device + referrer)", GREEN),
                ("orders / catalog\n(derived tables)", GREEN),
                ("funnel_metrics/\nattribution/ anomalies/", BLUE),
            ],
        },
        {
            "label": "CORE PIPELINE",
            "x": Inches(3.42),
            "w": Inches(3.1),
            "color": BLUE,
            "items": [
                ("etl_enrichment.py\nCRC32 hash · Parquet write · 229 lines", GREEN),
                ("benchmark_job.py\nOrchestrator · Sessionization · Funnel · 259 lines", BLUE),
                ("attribution.py\nLast-touch join · row_number() · 285 lines", RED),
                ("anomaly_detection.py\n7-day z-score · SPIKE/DROP text · 279 lines", PURPLE),
            ],
        },
        {
            "label": "OPTIMIZATIONS",
            "x": Inches(6.86),
            "w": Inches(2.9),
            "color": ORANGE,
            "items": [
                ("broadcast_join_opt.py\nShuffle vs broadcast join · 213 lines", ORANGE),
                ("partition_analysis.py\nDate-partitioned Parquet · 177 lines", ORANGE),
                ("skew_analysis.py\nHot-key detection · salting · 315 lines", ORANGE),
            ],
        },
        {
            "label": "CLOUD / INFRA",
            "x": Inches(10.10),
            "w": Inches(2.95),
            "color": TEAL,
            "items": [
                ("cloud_pipeline.py\nFull Dataproc pipeline · 554 lines", TEAL),
                ("cloud_submit.sh\nParameterized job submission", TEAL),
                ("demo.sh\nLive demo entry point", TEAL),
                ("streaming_job.py\nStructured streaming demo", BLUE),
            ],
        },
    ]

    col_xs = [c["x"] for c in columns]
    col_ws = [c["w"] for c in columns]

    for col in columns:
        cx, cw = col["x"], col["w"]
        # Column label
        lbl_box = box(slide, cx, COL_Y, cw, COL_H_LBL, fill=col["color"], rounded=True)
        txt(slide, col["label"], cx, COL_Y, cw, COL_H_LBL,
            size=Pt(11), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        # Items
        iy = COL_Y + COL_H_LBL + Inches(0.14)
        for (item_lbl, item_color) in col["items"]:
            b = box(slide, cx, iy, cw, ROW_H, fill=item_color, rounded=True)
            lines = item_lbl.split("\n")
            # bold first line, normal second
            txt(slide, lines[0], cx, iy + Inches(0.05), cw, Inches(0.34),
                size=Pt(11), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
            if len(lines) > 1:
                txt(slide, lines[1], cx, iy + Inches(0.38), cw, Inches(0.42),
                    size=Pt(9.5), color=OFFWHITE, align=PP_ALIGN.CENTER)
            iy += ROW_H + ROW_GAP

    # ── Horizontal flow arrows (at mid-height between col label and items) ──
    arr_y = COL_Y + COL_H_LBL + Inches(0.14) + ROW_H * 1.5 + ROW_GAP
    for i in range(len(columns) - 1):
        ax2 = col_xs[i] + col_ws[i] + Inches(0.02)
        aw  = col_xs[i+1] - ax2 - Inches(0.02)
        box(slide, ax2, arr_y, aw, Inches(0.04), fill=LGRAY)
        # arrowhead (small triangle drawn as tiny box)
        box(slide, col_xs[i+1] - Inches(0.06), arr_y - Inches(0.06),
            Inches(0.08), Inches(0.16), fill=LGRAY, rounded=True)

    # ── Dependency note ──────────────────────────────────────────────────
    notes = [
        "Raw → Processed", "Reads Parquet", "Uses Parquet", "Runs on GCP"
    ]
    for i, (col, note) in enumerate(zip(columns, notes)):
        txt(slide, note, col["x"], Inches(6.9), col["w"], Inches(0.32),
            size=Pt(10), color=GRAY, align=PP_ALIGN.CENTER, italic=True)

    # ── Total LOC banner ─────────────────────────────────────────────────
    box(slide, Inches(0.28), Inches(7.12), Inches(12.76), Inches(0.28), fill=NAVY)
    txt(slide,
        "Total codebase: ~4,080 lines of Python  ·  17 scripts  ·  "
        "Local (PySpark local[*])  +  Cloud (GCP Dataproc)",
        Inches(0.36), Inches(7.12), Inches(12.7), Inches(0.28),
        size=Pt(11), bold=False, color=LGRAY, align=PP_ALIGN.CENTER)
    return slide


def slide_pipeline(prs):
    """Slide 5 — Pipeline: 5 Stages."""
    slide = blank_slide(prs)
    set_bg(slide, OFFWHITE)
    header(slide, "Pipeline Deep Dive — 5 Stages", "Slide 4 / 7")

    stages = [
        {
            "num": "1", "name": "ETL &\nEnrichment", "color": GREEN,
            "points": [
                "Read 42M-row CSV via spark.read.csv",
                "CRC32(user_session) → device, referrer",
                "Filter purchases → orders table",
                "Write Parquet to GCS",
            ],
            "kpi": "224s write\n(5w / 10GB)",
        },
        {
            "num": "2", "name": "Session-\nization", "color": BLUE,
            "points": [
                "lag(event_time) per user_id window",
                "Gap > 30 min → new session flag",
                "cumsum() → session index",
                "session_id = user_id + '_' + idx",
            ],
            "kpi": "185s write\n(5w / 10GB)",
        },
        {
            "num": "3", "name": "Funnel\nAnalysis", "color": ORANGE,
            "points": [
                "view → cart → purchase events",
                "Group by category / device / referrer",
                "Daily and weekly granularity",
                "28,460 funnel rows (1GB run)",
            ],
            "kpi": "28,460\nfunnel rows",
        },
        {
            "num": "4", "name": "Attribution\n(Last-Touch)", "color": RED,
            "points": [
                "Join orders ⊲⊳ sessions on user_id",
                "Filter: referrer ≠ 'direct'",
                "24-hour lookback window",
                "row_number() → most recent ref.",
            ],
            "kpi": "14.2s join\n(5w / 1GB)",
        },
        {
            "num": "5", "name": "Anomaly\nDetection", "color": PURPLE,
            "points": [
                "Daily purchase counts (system + per-cat)",
                "7-day rolling avg + stddev window",
                "Flag |z-score| > 2 as SPIKE / DROP",
                "8 system / 305 category anomalies",
            ],
            "kpi": "8 sys.\n305 cat.",
        },
    ]

    CW   = Inches(2.45)
    CH   = Inches(6.02)
    CX0  = Inches(0.28)
    CY   = Inches(0.76)
    CGAP = Inches(0.13)

    for i, s in enumerate(stages):
        cx = CX0 + i * (CW + CGAP)

        # Card
        box(slide, cx, CY, CW, CH, fill=WHITE, line=LGRAY, lw=Pt(1))

        # Header
        box(slide, cx, CY, CW, Inches(0.74), fill=s["color"])
        txt(slide, f"Stage {s['num']}", cx, CY + Inches(0.04), CW, Inches(0.3),
            size=Pt(10), color=OFFWHITE, align=PP_ALIGN.CENTER)
        txt(slide, s["name"], cx, CY + Inches(0.3), CW, Inches(0.48),
            size=Pt(13.5), bold=True, color=WHITE, align=PP_ALIGN.CENTER)

        # Bullet points
        by = CY + Inches(0.84)
        for pt in s["points"]:
            box(slide, cx + Inches(0.2), by + Inches(0.11),
                Inches(0.1), Inches(0.1), fill=s["color"], rounded=True)
            txt(slide, pt, cx + Inches(0.38), by, CW - Inches(0.5), Inches(0.58),
                size=Pt(10.5), color=NAVY, wrap=True)
            by += Inches(0.82)

        # KPI badge at bottom
        kpi_y = CY + CH - Inches(0.62)
        box(slide, cx + Inches(0.18), kpi_y, CW - Inches(0.36), Inches(0.52),
            fill=s["color"], rounded=True)
        txt(slide, s["kpi"], cx + Inches(0.18), kpi_y, CW - Inches(0.36), Inches(0.52),
            size=Pt(10.5), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    return slide


def slide_optimizations(prs):
    """Slide 6 — Performance Optimizations."""
    slide = blank_slide(prs)
    set_bg(slide, OFFWHITE)
    header(slide, "Performance Optimizations", "Slide 5 / 7")

    opt_chart = chart_optimizations()
    img(slide, opt_chart, Inches(0.28), Inches(0.76), Inches(12.76), Inches(4.5))

    # Summary cards
    cards = [
        ("Broadcast Join", "1.58×  faster\n11.33s → 7.15s",
         "Orders (153K rows) broadcast to all workers,\neliminating cross-node shuffle entirely.", GREEN),
        ("Partition Pruning", "3.11×  single-day\n1.73×  7-day range",
         "Date-partitioned Parquet lets Spark skip entire\ndate directories — pure I/O reduction.", ORANGE),
        ("Skew Analysis", "Salting was 2.28× SLOWER\n34.6s → 79.1s",
         "Top-1 user = 0.02% of events — no real skew.\nSalting crossJoin overhead outweighed benefit.", BLUE),
    ]
    cw, ch = Inches(4.08), Inches(1.1)
    cx0 = Inches(0.28)
    cy = Inches(5.38)
    for i, (title, metric, desc, color) in enumerate(cards):
        cx = cx0 + i * (cw + Inches(0.16))
        box(slide, cx, cy, cw, ch, fill=color, rounded=True)
        txt(slide, title, cx, cy + Inches(0.04), cw, Inches(0.32),
            size=Pt(12), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        txt(slide, metric, cx, cy + Inches(0.34), cw, Inches(0.38),
            size=Pt(13), bold=True, color=GOLD, align=PP_ALIGN.CENTER)
        txt(slide, desc, cx, cy + Inches(0.7), cw, Inches(0.38),
            size=Pt(9.5), color=OFFWHITE, align=PP_ALIGN.CENTER)
    return slide


def slide_results(prs):
    """Slide 7 — Evaluation Results."""
    slide = blank_slide(prs)
    set_bg(slide, OFFWHITE)
    header(slide, "Evaluation Results — Scaling Benchmarks", "Slide 6 / 7")

    # Left: execution time chart
    time_chart = chart_scaling_time()
    img(slide, time_chart, Inches(0.28), Inches(0.76), Inches(7.7), Inches(4.1))

    # Right: speedup chart
    spdup_chart = chart_speedup()
    img(slide, spdup_chart, Inches(8.14), Inches(0.76), Inches(5.0), Inches(4.1))

    # Benchmark table
    txt(slide, "Full Benchmark Summary  (Total Pipeline Time in Seconds)",
        Inches(0.28), Inches(4.95), Inches(12.76), Inches(0.34),
        size=Pt(12), bold=True, color=NAVY)

    col_headers = ["Config", "1 GB  (~8.5M rows)", "5 GB  (~42.5M rows)",
                   "10 GB  (~85M rows)", "Throughput  (10 GB)"]
    cw_list = [Inches(1.5), Inches(2.6), Inches(2.6), Inches(2.6), Inches(3.02)]
    rx, ry, rh = Inches(0.28), Inches(5.33), Inches(0.35)

    # Header row
    tx = rx
    for hdr, cw in zip(col_headers, cw_list):
        box(slide, tx, ry, cw, rh, fill=NAVY)
        txt(slide, hdr, tx, ry, cw, rh, size=Pt(10), bold=True,
            color=WHITE, align=PP_ALIGN.CENTER)
        tx += cw

    rows = [
        ("2 Workers", "673.6 s", "1,021.8 s", "1,378.9 s", "61,569 rows/s"),
        ("4 Workers", "282.4 s",   "569.0 s",   "799.8 s", "106,148 rows/s"),
        ("5 Workers", "256.9 s",   "409.7 s",   "674.0 s", "125,961 rows/s"),
    ]
    alt = [RGBColor(0xE8, 0xF0, 0xFA), WHITE, RGBColor(0xE8, 0xF0, 0xFA)]
    for k, (row, fill) in enumerate(zip(rows, alt)):
        tx = rx
        ry2 = ry + rh + k * rh
        for j, (cell, cw) in enumerate(zip(row, cw_list)):
            cfill = NAVY if j == 0 else fill
            box(slide, tx, ry2, cw, rh, fill=cfill, line=LGRAY, lw=Pt(0.5))
            ctxt  = WHITE if j == 0 else NAVY
            txt(slide, cell, tx, ry2, cw, rh, size=Pt(10.5),
                bold=(j == 0), color=ctxt, align=PP_ALIGN.CENTER)
            tx += cw

    # Highlight 5-worker gain
    txt(slide, "★  5 workers vs. 2 workers on 10 GB:  2.05× speedup  (674s → 1379s baseline)",
        Inches(0.28), Inches(7.0), Inches(12.76), Inches(0.32),
        size=Pt(12), bold=True, color=BLUE, align=PP_ALIGN.CENTER)
    return slide


def slide_lessons(prs):
    """Slide 8 — Lessons Learned & Conclusion."""
    slide = blank_slide(prs)
    set_bg(slide, OFFWHITE)
    header(slide, "Lessons Learned & Conclusion", "Slide 7 / 7")

    lessons = [
        {
            "id": "L1",
            "title": "ETL I/O Dominates at Scale",
            "body": ("At 10 GB, ETL Parquet write consumed 477 s out of 1,379 s total (35% of "
                     "pipeline runtime). Distributed I/O is the primary bottleneck — not computation. "
                     "Profile write strategies before tuning CPU-bound stages."),
            "evidence": "ETL write: 477 s  /  Total: 1,379 s  =  35%",
            "color": GREEN,
        },
        {
            "id": "L2",
            "title": "Salting Hurts on Low-Skew Data",
            "body": ("Applying 20-bucket salting to the attribution join was 2.28× slower "
                     "(34.6 s → 79.1 s). The dataset had no meaningful hot key (top-1 user = 0.02%). "
                     "Always measure skew before applying mitigation."),
            "evidence": "No salting: 34.6 s  →  With salting: 79.1 s  (0.44× speedup)",
            "color": RED,
        },
        {
            "id": "L3",
            "title": "Broadcast Join Eliminates Network Shuffle",
            "body": ("Broadcasting the smaller orders table (153 K rows) gave a 1.58× speedup "
                     "by avoiding cross-node data transfer. The benefit grows with cluster size "
                     "as more network hops are eliminated."),
            "evidence": "Shuffle join: 11.33 s  →  Broadcast: 7.15 s  (1.58×)",
            "color": BLUE,
        },
        {
            "id": "L4",
            "title": "Storage Layout Beats Raw Compute Scaling",
            "body": ("Date-partitioned Parquet gave a 3.11× speedup on a single-day query — "
                     "exceeding the 2.05× gained by growing the cluster from 2 → 5 workers. "
                     "Data layout decisions often have larger impact than adding hardware."),
            "evidence": "Flat query: 4.25 s  →  Partitioned: 1.37 s  (3.11×)",
            "color": ORANGE,
        },
    ]

    LW, LH = Inches(6.22), Inches(2.7)
    xs = [Inches(0.28), Inches(6.82)]
    ys = [Inches(0.80), Inches(3.65)]

    for i, l in enumerate(lessons):
        lx = xs[i % 2]
        ly = ys[i // 2]
        box(slide, lx, ly, LW, LH, fill=WHITE, line=l["color"], lw=Pt(2.2))
        # Colored side bar
        box(slide, lx, ly, Inches(0.52), LH, fill=l["color"])
        txt(slide, l["id"], lx, ly + LH/2 - Inches(0.22),
            Inches(0.52), Inches(0.44),
            size=Pt(14), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        # Content
        txt(slide, l["title"], lx + Inches(0.62), ly + Inches(0.1),
            LW - Inches(0.76), Inches(0.42),
            size=Pt(14), bold=True, color=l["color"])
        txt(slide, l["body"], lx + Inches(0.62), ly + Inches(0.52),
            LW - Inches(0.76), Inches(1.4),
            size=Pt(11), color=NAVY, wrap=True)
        # Evidence strip
        ev_y = ly + LH - Inches(0.52)
        box(slide, lx + Inches(0.62), ev_y, LW - Inches(0.76), Inches(0.4),
            fill=l["color"], rounded=True)
        txt(slide, l["evidence"], lx + Inches(0.62), ev_y, LW - Inches(0.76), Inches(0.4),
            size=Pt(10), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    return slide


# ─────────────────────────────────────────────────────────────────────────────
#  MAIN
# ─────────────────────────────────────────────────────────────────────────────

def main():
    prs = new_prs()

    steps = [
        ("1/8  Title",                   slide_title),
        ("2/8  Problem Statement",        slide_problem),
        ("3/8  System Architecture",      slide_architecture),
        ("4/8  Implementation Structure", slide_implementation),
        ("5/8  Pipeline Stages",          slide_pipeline),
        ("6/8  Optimizations",            slide_optimizations),
        ("7/8  Scaling Results",          slide_results),
        ("8/8  Lessons Learned",          slide_lessons),
    ]

    for label, fn in steps:
        print(f"  Building slide {label} ...", flush=True)
        fn(prs)

    prs.save(OUT_PATH)
    print(f"\nSaved → {OUT_PATH}")
    print(f"Slides : {len(prs.slides)}")


if __name__ == "__main__":
    main()
